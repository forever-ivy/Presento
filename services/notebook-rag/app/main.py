import base64
import csv
import io
import os
import re
import zipfile
from typing import Any, Literal

from fastapi import Depends, FastAPI, Header, HTTPException
from pydantic import BaseModel, Field


app = FastAPI(title="Presento Notebook RAG Sidecar", version="0.1.0")


class ParseFileRequest(BaseModel):
    fileId: str
    fileName: str
    fileKind: str
    mimeType: str | None = None
    storagePath: str | None = None
    storageKey: str | None = None
    signedUrl: str | None = None
    contentBase64: str | None = None


class ExplainChunk(BaseModel):
    content: str
    metadata: dict[str, Any] = Field(default_factory=dict)


class ExplainFileRequest(BaseModel):
    fileId: str
    fileName: str
    mode: Literal["quick", "mastery"] = "quick"
    chunks: list[ExplainChunk] = Field(default_factory=list)
    question: str | None = None


def require_api_key(x_api_key: str | None = Header(default=None)) -> None:
    expected = os.getenv("NOTEBOOK_RAG_API_KEY")
    if expected and x_api_key != expected:
        raise HTTPException(status_code=401, detail="Invalid sidecar API key.")


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/parse-file", dependencies=[Depends(require_api_key)])
def parse_file(payload: ParseFileRequest) -> dict[str, Any]:
    data = decode_payload(payload)
    text, extras = extract_text(payload.fileName, payload.fileKind, data)
    chunks = chunk_text(text, payload.fileName, extras.get("chunkMetadata", {}))
    outline = infer_outline(text)
    return {
        "source": {
            "title": title_from_filename(payload.fileName),
            "summary": summarize_text(text, payload.fileName),
            "fileKind": payload.fileKind,
            "metadata": {
                "mimeType": payload.mimeType,
                "storagePath": payload.storagePath,
                "storageKey": payload.storageKey,
                **extras.get("sourceMetadata", {}),
            },
        },
        "chunks": chunks,
        "preview": {
            "text": "\n".join(text.splitlines()[:24]),
            "outline": outline,
        },
        "slides": extras.get("slides", []),
        "tables": extras.get("tables", []),
        "codeTree": extras.get("codeTree", []),
        "citations": [citation_from_chunk(payload.fileName, chunk) for chunk in chunks[:8]],
    }


@app.post("/explain-file", dependencies=[Depends(require_api_key)])
def explain_file(payload: ExplainFileRequest) -> dict[str, Any]:
    return explain_from_chunks(payload)


@app.post("/chat-file", dependencies=[Depends(require_api_key)])
def chat_file(payload: ExplainFileRequest) -> dict[str, Any]:
    if not payload.question:
        raise HTTPException(status_code=400, detail="Missing question.")
    return explain_from_chunks(payload)


def decode_payload(payload: ParseFileRequest) -> bytes:
    if payload.contentBase64:
        return base64.b64decode(payload.contentBase64)
    raise HTTPException(status_code=400, detail="Only inline contentBase64 is enabled in the MVP sidecar.")


def extract_text(file_name: str, file_kind: str, data: bytes) -> tuple[str, dict[str, Any]]:
    lower = file_name.lower()
    if lower.endswith(".pdf"):
        return extract_pdf(data), {"chunkMetadata": {"parser": "pypdf"}}
    if lower.endswith(".docx"):
        return extract_docx(data), {"chunkMetadata": {"parser": "python-docx"}}
    if lower.endswith((".pptx", ".ppt")):
        text, slides = extract_pptx(data)
        return text, {"slides": slides, "chunkMetadata": {"parser": "python-pptx"}}
    if lower.endswith((".xlsx", ".xlsm")):
        text, tables = extract_xlsx(data)
        return text, {"tables": tables, "chunkMetadata": {"parser": "openpyxl"}}
    if lower.endswith(".zip") or file_kind == "code":
        text, tree = extract_code_zip(data)
        return text, {"codeTree": tree, "chunkMetadata": {"parser": "zip-repomix-lite"}}
    if lower.endswith(".csv"):
        return extract_csv(data), {"chunkMetadata": {"parser": "csv"}}
    return decode_text(data), {"chunkMetadata": {"parser": "plain-text"}}


def extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            pages.append(f"[page {index}]\n{page.extract_text() or ''}")
        return "\n\n".join(pages)
    except Exception as exc:
        return f"PDF parsing failed: {exc}"


def extract_docx(data: bytes) -> str:
    try:
        from docx import Document

        document = Document(io.BytesIO(data))
        return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    except Exception as exc:
        return f"DOCX parsing failed: {exc}"


def extract_pptx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    try:
        from pptx import Presentation

        deck = Presentation(io.BytesIO(data))
        slides = []
        for index, slide in enumerate(deck.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
            title = text_runs[0] if text_runs else f"Slide {index}"
            slide_text = "\n".join(text_runs)
            slides.append({"page": index, "title": title[:120], "text": slide_text, "metadata": {}})
        return "\n\n".join(f"[slide {item['page']}]\n{item['text']}" for item in slides), slides
    except Exception as exc:
        return f"PPTX parsing failed: {exc}", []


def extract_xlsx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        text_blocks = []
        tables = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(max_row=50, values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    rows.append(values)
            if not rows:
                continue
            headers = rows[0]
            tables.append({"sheet": sheet.title, "headers": headers, "rows": rows[1:21], "metadata": {}})
            text_blocks.append(f"[sheet {sheet.title}]\n" + "\n".join("\t".join(row) for row in rows[:50]))
        return "\n\n".join(text_blocks), tables
    except Exception as exc:
        return f"XLSX parsing failed: {exc}", []


def extract_csv(data: bytes) -> str:
    decoded = decode_text(data)
    rows = list(csv.reader(io.StringIO(decoded)))
    return "\n".join("\t".join(row) for row in rows[:200])


def extract_code_zip(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    tree = []
    blocks = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        for name in archive.namelist():
            if name.endswith("/") or is_binary_path(name):
                continue
            raw = archive.read(name)
            text = decode_text(raw)
            line_count = len(text.splitlines())
            tree.append({
                "path": name,
                "language": language_for_path(name),
                "summary": f"{line_count} lines",
                "lineCount": line_count,
            })
            blocks.append(f"--- {name} ---\n{text[:12000]}")
    return "\n\n".join(blocks), tree


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def chunk_text(text: str, file_name: str, metadata: dict[str, Any]) -> list[dict[str, Any]]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    chunks = []
    current: list[str] = []
    start_line = 1
    for index, line in enumerate(lines, start=1):
        if current and sum(len(item) for item in current) + len(line) > 1200:
            chunks.append(make_chunk(file_name, current, start_line, index - 1, metadata))
            current = []
            start_line = index
        current.append(line)
    if current:
        chunks.append(make_chunk(file_name, current, start_line, len(lines), metadata))
    return chunks


def make_chunk(file_name: str, lines: list[str], start_line: int, end_line: int, metadata: dict[str, Any]) -> dict[str, Any]:
    content = "\n".join(lines)
    page = infer_marker_number(content, "page")
    slide = infer_marker_number(content, "slide")
    return {
        "content": content,
        "source": file_name,
        "metadata": {
            **metadata,
            "fileName": file_name,
            "lineStart": start_line,
            "lineEnd": end_line,
            **({"page": page} if page else {}),
            **({"slide": slide} if slide else {}),
        },
    }


def explain_from_chunks(payload: ExplainFileRequest) -> dict[str, Any]:
    if not payload.chunks:
        return {
            "summary": "依据不足：当前文件没有可检索片段。",
            "outline": [],
            "answer": "依据不足：我不能在没有文件证据的情况下编造答案。",
            "citations": [],
        }

    chunks = payload.chunks[: 4 if payload.mode == "quick" else 10]
    outline = [first_line(chunk.content) for chunk in chunks if first_line(chunk.content)]
    citations = [citation_from_chunk(payload.fileName, chunk.model_dump()) for chunk in chunks]
    if payload.question:
        answer = f"围绕“{payload.question}”，可以依据文件这样回答：{'; '.join(outline[:4])}。"
        summary = f"{payload.fileName} 追问回答"
    elif payload.mode == "quick":
        summary = f"{payload.fileName} 的速通重点是：{', '.join(outline[:3])}。"
        answer = "\n".join([
            f"一句话：{summary}",
            f"3-5 个要点：{'; '.join(outline[:5])}",
            "高危追问：技术选择依据、个人贡献、结果证据。",
            "30 秒框架：目标 -> 方法 -> 证据 -> 我的贡献。",
        ])
    else:
        summary = f"{payload.fileName} 的精通讲解覆盖 {len(outline)} 个证据片段。"
        answer = "\n".join([
            summary,
            "分段讲解：",
            *[f"{index + 1}. {item}" for index, item in enumerate(outline)],
            "自测题：请解释核心目标、关键证据、个人贡献和可能追问。",
        ])
    return {
        "summary": summary,
        "outline": outline,
        "answer": answer,
        "citations": citations,
        "followUps": ["把它压缩成 30 秒回答", "列出最可能被追问的 3 个点"],
        "quiz": ["核心证据是什么？", "你负责了哪一部分？"] if payload.mode == "mastery" else [],
        "weaknessCandidates": ["证据引用不熟", "个人贡献边界不清"],
        "metadata": {"engine": "sidecar-deterministic"},
    }


def citation_from_chunk(file_name: str, chunk: dict[str, Any]) -> dict[str, Any]:
    metadata = chunk.get("metadata") or {}
    return {
        "fileName": metadata.get("fileName", file_name),
        "page": metadata.get("page"),
        "slide": metadata.get("slide"),
        "sheet": metadata.get("sheet"),
        "cellRange": metadata.get("cellRange"),
        "codePath": metadata.get("codePath"),
        "lineStart": metadata.get("lineStart"),
        "lineEnd": metadata.get("lineEnd"),
        "text": first_line(chunk.get("content") or ""),
    }


def infer_outline(text: str) -> list[str]:
    headings = []
    for line in text.splitlines():
        cleaned = re.sub(r"^#+\s*", "", line).strip()
        if cleaned and (line.startswith("#") or len(cleaned) < 60):
            headings.append(cleaned)
        if len(headings) >= 8:
            break
    return headings


def summarize_text(text: str, file_name: str) -> str:
    line_count = len([line for line in text.splitlines() if line.strip()])
    return f"{file_name} 已解析为 {line_count} 行可学习内容。"


def title_from_filename(file_name: str) -> str:
    return re.sub(r"\.[^.]+$", "", file_name).strip() or file_name


def first_line(text: str) -> str:
    return next((line.strip() for line in text.splitlines() if line.strip()), "")


def infer_marker_number(text: str, marker: str) -> int | None:
    match = re.search(rf"\[{marker}\s+(\d+)\]", text, flags=re.IGNORECASE)
    return int(match.group(1)) if match else None


def language_for_path(path: str) -> str:
    extension = path.rsplit(".", 1)[-1].lower() if "." in path else ""
    return {
        "ts": "typescript",
        "tsx": "typescriptreact",
        "js": "javascript",
        "jsx": "javascriptreact",
        "py": "python",
        "sql": "sql",
        "md": "markdown",
        "java": "java",
        "kt": "kotlin",
    }.get(extension, extension or "text")


def is_binary_path(path: str) -> bool:
    return path.lower().endswith((
        ".png", ".jpg", ".jpeg", ".gif", ".webp", ".pdf", ".pptx", ".docx", ".xlsx",
        ".zip", ".ico", ".woff", ".woff2", ".ttf",
    ))
